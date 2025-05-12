import os
from typing import List, Dict, Any
import google.generativeai as genai
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_huggingface import HuggingFaceEmbeddings  # updated import
from langchain_community.vectorstores import FAISS
import PyPDF2
from docx import Document
import pandas as pd
import magic
import tiktoken
import json
import xml.etree.ElementTree as ET
import yaml
import csv
import xlrd
import openpyxl
import pptx
import html
import re

class DocumentProcessor:
    def __init__(self, api_key: str):
        self.api_key = api_key
        genai.configure(api_key=api_key)
        self.gemini_model = genai.GenerativeModel('models/gemini-2.0-flash-lite')
        self.embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
        
    def read_file(self, file_path: str) -> str:
        """Read different file types and return their content as text."""
        mime = magic.Magic(mime=True)
        file_type = mime.from_file(file_path)
        file_extension = os.path.splitext(file_path)[1].lower()
        
        # Handle different file types
        if file_type == 'application/pdf':
            return self._read_pdf(file_path)
        elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            return self._read_docx(file_path)
        elif file_type == 'text/plain':
            return self._read_txt(file_path)
        elif file_type == 'text/csv':
            return self._read_csv(file_path)
        elif file_type == 'application/json':
            return self._read_json(file_path)
        elif file_type == 'text/xml' or file_extension == '.xml':
            return self._read_xml(file_path)
        elif file_type == 'text/yaml' or file_extension in ['.yml', '.yaml']:
            return self._read_yaml(file_path)
        elif file_type == 'application/vnd.ms-excel' or file_extension == '.xls':
            return self._read_xls(file_path)
        elif file_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet' or file_extension == '.xlsx':
            return self._read_xlsx(file_path)
        elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' or file_extension == '.pptx':
            return self._read_pptx(file_path)
        elif file_type == 'text/html' or file_extension == '.html':
            return self._read_html(file_path)
        else:
            # Try to read as text if file type is unknown
            try:
                return self._read_txt(file_path)
            except:
                raise ValueError(f"Unsupported file type: {file_type}")
    
    def _read_pdf(self, file_path: str) -> str:
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text
    
    def _read_docx(self, file_path: str) -> str:
        doc = Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])
    
    def _read_txt(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    
    def _read_csv(self, file_path: str) -> str:
        df = pd.read_csv(file_path)
        return df.to_string()
    
    def _read_json(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            data = json.load(file)
            return json.dumps(data, indent=2)
    
    def _read_xml(self, file_path: str) -> str:
        tree = ET.parse(file_path)
        root = tree.getroot()
        return ET.tostring(root, encoding='unicode', method='xml')
    
    def _read_yaml(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            data = yaml.safe_load(file)
            return yaml.dump(data, default_flow_style=False)
    
    def _read_xls(self, file_path: str) -> str:
        workbook = xlrd.open_workbook(file_path)
        text = ""
        for sheet in workbook.sheets():
            text += f"Sheet: {sheet.name}\n"
            for row in range(sheet.nrows):
                text += "\t".join(str(sheet.cell_value(row, col)) for col in range(sheet.ncols)) + "\n"
        return text
    
    def _read_xlsx(self, file_path: str) -> str:
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        text = ""
        for sheet in workbook.sheetnames:
            text += f"Sheet: {sheet}\n"
            ws = workbook[sheet]
            for row in ws.rows:
                text += "\t".join(str(cell.value) for cell in row) + "\n"
        return text
    
    def _read_pptx(self, file_path: str) -> str:
        prs = pptx.Presentation(file_path)
        text = ""
        for slide in prs.slides:
            text += f"Slide {slide.slide_number}\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    def _read_html(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as file:
            html_content = file.read()
            # Remove HTML tags
            text = re.sub('<[^<]+?>', '', html_content)
            # Decode HTML entities
            text = html.unescape(text)
            return text
    
    def process_documents(self, file_paths: List[str]) -> FAISS:
        """Process multiple documents and create a FAISS vector store."""
        all_texts = []
        for file_path in file_paths:
            try:
                text = self.read_file(file_path)
                all_texts.append(text)
            except Exception as e:
                print(f"Error processing {file_path}: {str(e)}")
        
        # Split texts into chunks
        chunks = []
        for text in all_texts:
            chunks.extend(self.text_splitter.split_text(text))
        
        # Create vector store
        vector_store = FAISS.from_texts(chunks, self.embeddings)
        return vector_store

    def ask_gemini(self, prompt: str, generation_config=None) -> str:
        """
        Get a response from Gemini using the Google Generative AI library,
        using a chat session for better context handling.
        """
        chat_session = self.gemini_model.start_chat(history=[])
        if generation_config is None:
            from google.generativeai.types import GenerationConfig
            generation_config = GenerationConfig(
                temperature=0.2,
                top_p=0.8,
                top_k=40
            )
        try:
            response = chat_session.send_message(prompt, generation_config=generation_config)
            return response.text
        except Exception as e:
            print(f"Error calling Gemini API: {str(e)}")
            return None